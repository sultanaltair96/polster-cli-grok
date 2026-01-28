#!/usr/bin/env python3
"""
Polster Presentation Generator
Creates a stunning PowerPoint presentation for data analytics meetup
Theme: White background with Polars Blue and Dagster Purple accents
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, Circle, Rectangle, FancyArrowPatch
import numpy as np
from io import BytesIO
import os

# Color Palette
COLORS = {
    "white": RGBColor(255, 255, 255),
    "bg_light": RGBColor(248, 250, 252),
    "text_dark": RGBColor(15, 23, 42),
    "text_gray": RGBColor(100, 116, 139),
    "polars_blue": RGBColor(93, 173, 226),
    "polars_blue_dark": RGBColor(40, 116, 166),
    "dagster_purple": RGBColor(124, 58, 237),
    "dagster_purple_dark": RGBColor(109, 40, 217),
    "bronze": RGBColor(205, 127, 50),
    "silver": RGBColor(192, 192, 192),
    "gold": RGBColor(255, 215, 0),
    "warning": RGBColor(251, 146, 60),
    "success": RGBColor(34, 197, 94),
}


class PolsterPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.setup_slide_master()

    def setup_slide_master(self):
        """Setup the slide master with consistent styling"""
        slide_master = self.prs.slide_master

        # Set background to white
        background = slide_master.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["white"]

    def add_text_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=18,
        font_color=None,
        bold=False,
        alignment=PP_ALIGN.LEFT,
        font_name="Inter",
    ):
        """Add a text box with consistent styling"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color or COLORS["text_dark"]
        p.alignment = alignment
        p.font.name = font_name

        return txBox

    def add_gradient_title(self, slide, text, left, top, width, height):
        """Add a title with gradient effect (simulated with blue to purple)"""
        # Add purple text
        txBox = self.add_text_box(
            slide,
            left,
            top,
            width,
            height,
            text,
            font_size=48,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.font.color.rgb = COLORS["dagster_purple"]
        return txBox

    def create_infographic_pipeline_chaos(self):
        """Create infographic showing data pipeline chaos"""
        fig, ax = plt.subplots(figsize=(10, 4))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 4)
        ax.axis("off")

        # Panel 1: Spaghetti Code
        rect1 = FancyBboxPatch(
            (0.2, 0.5),
            2.8,
            3,
            boxstyle="round,pad=0.1",
            facecolor="#FEF3C7",
            edgecolor="#F59E0B",
            linewidth=2,
        )
        ax.add_patch(rect1)
        ax.text(1.6, 3.2, "Spaghetti Code", ha="center", fontsize=12, fontweight="bold")
        # Draw tangled lines
        for i in range(8):
            x1, y1 = 0.5 + i * 0.3, 2.5 - i * 0.2
            x2, y2 = 2.0 + np.sin(i) * 0.5, 1.0 + np.cos(i) * 0.5
            ax.plot([x1, x2], [y1, y2], "o-", color="#DC2626", linewidth=1.5, alpha=0.6)

        # Panel 2: Silent Failures
        rect2 = FancyBboxPatch(
            (3.6, 0.5),
            2.8,
            3,
            boxstyle="round,pad=0.1",
            facecolor="#FEE2E2",
            edgecolor="#EF4444",
            linewidth=2,
        )
        ax.add_patch(rect2)
        ax.text(
            5.0, 3.2, "Silent Failures", ha="center", fontsize=12, fontweight="bold"
        )
        # Draw broken chain
        chain_x = [3.9, 4.4, 4.9, 5.4, 6.0]
        chain_y = [1.8, 1.8, 1.8, 1.8, 1.8]
        for i in range(len(chain_x) - 1):
            if i == 2:  # Broken link
                ax.plot(
                    [chain_x[i], chain_x[i] + 0.25],
                    [chain_y[i], chain_y[i] + 0.3],
                    "o-",
                    color="#DC2626",
                    linewidth=3,
                )
                ax.plot(
                    [chain_x[i] + 0.35, chain_x[i + 1]],
                    [chain_y[i] + 0.3, chain_y[i + 1]],
                    "o-",
                    color="#DC2626",
                    linewidth=3,
                )
            else:
                ax.plot(
                    [chain_x[i], chain_x[i + 1]],
                    [chain_y[i], chain_y[i + 1]],
                    "o-",
                    color="#7C3AED",
                    linewidth=3,
                )

        # Panel 3: Knowledge Silos
        rect3 = FancyBboxPatch(
            (7.0, 0.5),
            2.8,
            3,
            boxstyle="round,pad=0.1",
            facecolor="#E0E7FF",
            edgecolor="#6366F1",
            linewidth=2,
        )
        ax.add_patch(rect3)
        ax.text(
            8.4, 3.2, "Knowledge Silos", ha="center", fontsize=12, fontweight="bold"
        )
        # Draw isolated islands
        for i, (x, y) in enumerate([(7.4, 2.0), (8.4, 1.3), (9.4, 2.0)]):
            circle = Circle(
                (x, y), 0.3, facecolor="#818CF8", edgecolor="#4F46E5", linewidth=2
            )
            ax.add_patch(circle)
            ax.text(
                x,
                y,
                f"T{i + 1}",
                ha="center",
                va="center",
                fontsize=10,
                color="white",
                fontweight="bold",
            )

        # Stats
        ax.text(
            1.6, 0.2, "68% of projects fail", ha="center", fontsize=9, style="italic"
        )
        ax.text(5.0, 0.2, "40% time debugging", ha="center", fontsize=9, style="italic")
        ax.text(8.4, 0.2, "Isolated teams", ha="center", fontsize=9, style="italic")

        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        plt.close()
        return buf

    def create_medallion_pyramid(self):
        """Create the Medallion Architecture pyramid infographic"""
        fig, ax = plt.subplots(figsize=(8, 6))
        ax.set_xlim(0, 8)
        ax.set_ylim(0, 6)
        ax.axis("off")

        # Gold Layer (Top)
        gold_rect = FancyBboxPatch(
            (2.5, 4),
            3,
            1.2,
            boxstyle="round,pad=0.05",
            facecolor="#FFD700",
            edgecolor="#B8860B",
            linewidth=3,
        )
        ax.add_patch(gold_rect)
        ax.text(
            4, 4.6, "GOLD", ha="center", va="center", fontsize=16, fontweight="bold"
        )
        ax.text(4, 4.2, "Business Intelligence", ha="center", va="center", fontsize=9)
        ax.text(4, 3.95, "Analytics & ML", ha="center", va="center", fontsize=9)

        # Silver Layer (Middle)
        silver_rect = FancyBboxPatch(
            (1.5, 2.5),
            5,
            1.2,
            boxstyle="round,pad=0.05",
            facecolor="#C0C0C0",
            edgecolor="#808080",
            linewidth=3,
        )
        ax.add_patch(silver_rect)
        ax.text(
            4, 3.1, "SILVER", ha="center", va="center", fontsize=16, fontweight="bold"
        )
        ax.text(4, 2.7, "Cleaned & Validated", ha="center", va="center", fontsize=9)
        ax.text(4, 2.45, "Structured Data", ha="center", va="center", fontsize=9)

        # Bronze Layer (Bottom)
        bronze_rect = FancyBboxPatch(
            (0.5, 1),
            7,
            1.2,
            boxstyle="round,pad=0.05",
            facecolor="#CD7F32",
            edgecolor="#8B4513",
            linewidth=3,
        )
        ax.add_patch(bronze_rect)
        ax.text(
            4, 1.6, "BRONZE", ha="center", va="center", fontsize=16, fontweight="bold"
        )
        ax.text(4, 1.2, "Raw Data Ingestion", ha="center", va="center", fontsize=9)
        ax.text(
            4, 0.95, "APIs • Databases • Files", ha="center", va="center", fontsize=9
        )

        # Flow arrows
        arrow1 = FancyArrowPatch(
            (4, 3.9),
            (4, 3.7),
            arrowstyle="->",
            mutation_scale=30,
            linewidth=3,
            color="#7C3AED",
        )
        ax.add_patch(arrow1)
        arrow2 = FancyArrowPatch(
            (4, 2.4),
            (4, 2.2),
            arrowstyle="->",
            mutation_scale=30,
            linewidth=3,
            color="#2874A6",
        )
        ax.add_patch(arrow2)

        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        plt.close()
        return buf

    def create_code_comparison(self):
        """Create side-by-side code comparison"""
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

        # Without Polster
        ax1.set_xlim(0, 10)
        ax1.set_ylim(0, 10)
        ax1.axis("off")
        ax1.set_title(
            "WITHOUT POLSTER\n(Messy & Complex)", fontsize=14, fontweight="bold", pad=20
        )

        code_without = [
            "# imports, setup, config...",
            "from dagster import Definitions, Asset",
            "import os, sys, yaml",
            "",
            "# Manual asset registration",
            "defs = Definitions(",
            "    assets=[asset1, asset2...],",
            "    jobs=[...],",
            "    schedules=[...]",
            ")",
            "",
            "# Complex dependency mgmt",
            "# Manual testing setup",
            "# Boilerplate everywhere",
            "# 40 lines of setup code",
        ]

        y_pos = 9
        for line in code_without:
            color = (
                "#DC2626"
                if "complex" in line.lower()
                or "manual" in line.lower()
                or "boilerplate" in line.lower()
                else "#374151"
            )
            ax1.text(0.5, y_pos, line, fontsize=9, family="monospace", color=color)
            y_pos -= 0.7

        rect1 = FancyBboxPatch(
            (0.2, 0.5),
            9.6,
            9,
            boxstyle="round,pad=0.1",
            facecolor="#FEE2E2",
            edgecolor="#EF4444",
            linewidth=2,
        )
        ax1.add_patch(rect1)
        ax1.text(
            5,
            0.8,
            "❌ 40 lines of boilerplate",
            ha="center",
            fontsize=11,
            fontweight="bold",
            color="#DC2626",
        )

        # With Polster
        ax2.set_xlim(0, 10)
        ax2.set_ylim(0, 10)
        ax2.axis("off")
        ax2.set_title(
            "WITH POLSTER\n(Clean & Simple)", fontsize=14, fontweight="bold", pad=20
        )

        code_with = [
            "from dagster import asset",
            "import polars as pl",
            "",
            "@asset",
            "def bronze_sales():",
            "    return pl.read_csv('sales.csv')",
            "",
            "@asset",
            "def silver_sales(bronze_sales):",
            "    return bronze_sales.clean()",
            "",
            "@asset",
            "def gold_metrics(silver_sales):",
            "    return silver_sales.groupby('region').sum()",
        ]

        y_pos = 9
        for line in code_with:
            if line.startswith("@"):
                color = "#7C3AED"
                weight = "bold"
            elif line.startswith("def"):
                color = "#2874A6"
                weight = "bold"
            else:
                color = "#374151"
                weight = "normal"
            ax2.text(
                0.5,
                y_pos,
                line,
                fontsize=9,
                family="monospace",
                color=color,
                fontweight=weight,
            )
            y_pos -= 0.7

        rect2 = FancyBboxPatch(
            (0.2, 0.5),
            9.6,
            9,
            boxstyle="round,pad=0.1",
            facecolor="#D1FAE5",
            edgecolor="#10B981",
            linewidth=2,
        )
        ax2.add_patch(rect2)
        ax2.text(
            5,
            0.8,
            "✅ 80% less code",
            ha="center",
            fontsize=11,
            fontweight="bold",
            color="#059669",
        )

        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        plt.close()
        return buf

    def create_metrics_chart(self):
        """Create metrics comparison chart"""
        fig, ax = plt.subplots(figsize=(10, 5))

        metrics = ["Setup Time", "Bug Detection", "Onboarding", "Maintenance"]
        before = [100, 30, 100, 100]  # Normalized to 100
        after = [3, 100, 20, 25]  # Percentages

        x = np.arange(len(metrics))
        width = 0.35

        bars1 = ax.bar(
            x - width / 2,
            before,
            width,
            label="Before Polster",
            color="#94A3B8",
            edgecolor="#64748B",
            linewidth=1.5,
        )
        bars2 = ax.bar(
            x + width / 2,
            after,
            width,
            label="After Polster",
            color="#7C3AED",
            edgecolor="#6D28D9",
            linewidth=1.5,
        )

        ax.set_ylabel("Relative Score", fontsize=12)
        ax.set_title(
            "Impact Metrics: Before vs After Polster",
            fontsize=14,
            fontweight="bold",
            pad=20,
        )
        ax.set_xticks(x)
        ax.set_xticklabels(metrics, fontsize=11)
        ax.legend(fontsize=10)
        ax.set_ylim(0, 120)

        # Add value labels on bars
        for bars in [bars1, bars2]:
            for bar in bars:
                height = bar.get_height()
                ax.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height + 2,
                    f"{int(height)}%",
                    ha="center",
                    va="bottom",
                    fontsize=9,
                )

        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.grid(axis="y", alpha=0.3, linestyle="--")

        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        plt.close()
        return buf

    def create_timeline(self):
        """Create the 5-minute timeline infographic"""
        fig, ax = plt.subplots(figsize=(12, 3))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 3)
        ax.axis("off")

        # Timeline line
        ax.plot([1, 9], [1.5, 1.5], "-", color="#CBD5E1", linewidth=4)

        # Time markers
        times = [1, 3, 5, 7, 9]
        labels = ["0:00", "0:30", "1:00", "3:00", "5:00"]
        steps = [
            "Init\nProject",
            "Add\nBronze",
            "Add\nSilver",
            "Add\nGold",
            "Run\nPipeline",
        ]

        colors = ["#2874A6", "#CD7F32", "#C0C0C0", "#FFD700", "#7C3AED"]

        for i, (x, time, step, color) in enumerate(zip(times, labels, steps, colors)):
            # Circle marker
            circle = Circle(
                (x, 1.5), 0.25, facecolor=color, edgecolor="white", linewidth=3
            )
            ax.add_patch(circle)

            # Time label
            ax.text(x, 1.0, time, ha="center", fontsize=10, fontweight="bold")

            # Step label
            ax.text(x, 2.0, step, ha="center", fontsize=9)

        # Arrows between steps
        for i in range(len(times) - 1):
            ax.annotate(
                "",
                xy=(times[i + 1] - 0.3, 1.5),
                xytext=(times[i] + 0.3, 1.5),
                arrowprops=dict(arrowstyle="->", color="#64748B", lw=2),
            )

        ax.text(
            5,
            0.3,
            "From Zero to Production Pipeline in 5 Minutes",
            ha="center",
            fontsize=12,
            fontweight="bold",
            style="italic",
        )

        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        plt.close()
        return buf

    def create_title_slide(self):
        """Slide 1: Title Card"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Main title with gradient effect (using purple)
        title_box = self.add_text_box(
            slide,
            1,
            2.5,
            11.333,
            1.5,
            "Polster",
            font_size=72,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        title_box.text_frame.paragraphs[0].font.color.rgb = COLORS["dagster_purple"]

        # Subtitle
        self.add_text_box(
            slide,
            1,
            4,
            11.333,
            1,
            "Build Reliable Data Pipelines with Enforced Architecture",
            font_size=24,
            alignment=PP_ALIGN.CENTER,
            font_color=COLORS["text_gray"],
        )

        # Decorative pipeline graphic
        img_buf = BytesIO()
        fig, ax = plt.subplots(figsize=(10, 1.5))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 1.5)
        ax.axis("off")

        # Draw flowing pipeline
        x = np.linspace(0, 10, 100)
        y = 0.75 + 0.2 * np.sin(x * 0.8)
        ax.plot(x, y, linewidth=8, color="#2874A6", alpha=0.7)
        ax.plot(x, y, linewidth=5, color="#7C3AED", alpha=0.8)

        # Add data dots
        for i in range(0, 100, 20):
            ax.scatter(x[i], y[i], s=100, color="white", zorder=5)
            ax.scatter(x[i], y[i], s=50, color="#2874A6", zorder=6)

        plt.savefig(
            img_buf, format="png", dpi=150, bbox_inches="tight", facecolor="white"
        )
        img_buf.seek(0)
        plt.close()

        slide.shapes.add_picture(img_buf, Inches(1.667), Inches(5.5), width=Inches(10))

        # Presenter info
        self.add_text_box(
            slide,
            1,
            6.5,
            11.333,
            0.5,
            "Data Analytics Meetup • January 2026",
            font_size=16,
            alignment=PP_ALIGN.CENTER,
            font_color=COLORS["text_gray"],
        )

    def create_problem_slide(self):
        """Slide 2: The Problem"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Why Do Data Projects Fail?",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Add infographic
        img_buf = self.create_infographic_pipeline_chaos()
        slide.shapes.add_picture(
            img_buf, Inches(0.5), Inches(1.5), width=Inches(12.333)
        )

        # Key insight box
        box_left = Inches(2)
        box_top = Inches(6)
        box_width = Inches(9.333)
        box_height = Inches(1)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Light yellow
        shape.line.color.rgb = RGBColor(245, 158, 11)
        shape.line.width = Pt(2)

        self.add_text_box(
            slide,
            2.2,
            6.2,
            8.9,
            0.6,
            "💡 The root cause: Lack of enforced architecture leads to technical debt",
            font_size=18,
            alignment=PP_ALIGN.CENTER,
            font_color=COLORS["text_dark"],
        )

    def create_medallion_slide(self):
        """Slide 3: Medallion Architecture"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "The Solution: Medallion Architecture",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Add pyramid infographic
        img_buf = self.create_medallion_pyramid()
        slide.shapes.add_picture(img_buf, Inches(2), Inches(1.5), width=Inches(5))

        # Benefits list on right
        benefits = [
            "✓ Quality improves at each layer",
            "✓ Clear data lineage & debugging",
            "✓ Consistent structure across teams",
            "✓ Enforced best practices",
            "✓ Separate concerns by layer",
        ]

        y_pos = 2
        for benefit in benefits:
            self.add_text_box(slide, 7.5, y_pos, 5, 0.6, benefit, font_size=18)
            y_pos += 0.6

        # Key message box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.5), Inches(5), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(224, 231, 255)
        box.line.color.rgb = COLORS["dagster_purple"]

        self.add_text_box(
            slide,
            7.7,
            5.7,
            4.6,
            0.6,
            "Like refining ore:\nRaw → Purified → Valuable",
            font_size=16,
            alignment=PP_ALIGN.CENTER,
            bold=True,
        )

    def create_polster_intro_slide(self):
        """Slide 4: Meet Polster"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Meet Polster: Architecture as Code",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Before/After comparison text
        self.add_text_box(
            slide,
            0.5,
            1.5,
            5.5,
            0.5,
            "BEFORE:",
            font_size=20,
            bold=True,
            font_color=COLORS["warning"],
        )

        before_items = [
            "• Manual project setup (2-3 days)",
            "• Inconsistent structure",
            "• Complex dependency management",
            "• No enforced patterns",
            "• Production surprises",
        ]
        y_pos = 2.0
        for item in before_items:
            self.add_text_box(
                slide,
                0.5,
                y_pos,
                5.5,
                0.4,
                item,
                font_size=16,
                font_color=COLORS["text_gray"],
            )
            y_pos += 0.4

        self.add_text_box(
            slide,
            6.5,
            1.5,
            6,
            0.5,
            "AFTER:",
            font_size=20,
            bold=True,
            font_color=COLORS["success"],
        )

        after_items = [
            "• One command setup (5 minutes)",
            "• Enforced Medallion layers",
            "• Automatic dependencies",
            "• Built-in best practices",
            "• Catch errors early",
        ]
        y_pos = 2.0
        for item in after_items:
            self.add_text_box(
                slide,
                6.5,
                y_pos,
                6,
                0.4,
                item,
                font_size=16,
                font_color=COLORS["text_gray"],
            )
            y_pos += 0.4

        # Command highlight box
        cmd_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2),
            Inches(5),
            Inches(9.333),
            Inches(1.5),
        )
        cmd_box.fill.solid()
        cmd_box.fill.fore_color.rgb = RGBColor(15, 23, 42)

        self.add_text_box(
            slide,
            2.2,
            5.2,
            9,
            0.4,
            "$ polster init my_pipeline",
            font_size=24,
            bold=True,
            font_color=COLORS["white"],
        )

        self.add_text_box(
            slide,
            2.2,
            5.7,
            9,
            0.6,
            "✓ Bronze layer created  ✓ Silver layer created  ✓ Gold layer created  ✓ Ready to build",
            font_size=14,
            font_color=COLORS["success"],
        )

    def create_demo_slide(self):
        """Slide 5: Live Demo Timeline"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "5 Minutes to Production",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Add timeline
        img_buf = self.create_timeline()
        slide.shapes.add_picture(
            img_buf, Inches(0.5), Inches(1.5), width=Inches(12.333)
        )

        # Step descriptions
        steps = [
            ("1. Initialize", "polster init my_project"),
            (
                "2. Add Bronze Asset",
                "polster add-asset --layer bronze --name raw_sales",
            ),
            (
                "3. Add Silver Asset",
                "polster add-asset --layer silver --name clean_sales --dependencies raw_sales",
            ),
            (
                "4. Add Gold Asset",
                "polster add-asset --layer gold --name sales_metrics --dependencies clean_sales",
            ),
            ("5. Run Pipeline", "python run_polster.py --ui"),
        ]

        y_pos = 4.5
        for title, cmd in steps:
            self.add_text_box(
                slide,
                0.5,
                y_pos,
                3,
                0.4,
                title,
                font_size=14,
                bold=True,
                font_color=COLORS["dagster_purple"],
            )
            self.add_text_box(
                slide,
                3.5,
                y_pos,
                9,
                0.4,
                f"$ {cmd}",
                font_size=12,
                font_color=COLORS["text_gray"],
                font_name="Courier New",
            )
            y_pos += 0.5

    def create_code_slide(self):
        """Slide 6: Code Comparison"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.3,
            12.333,
            0.8,
            "Less Boilerplate, More Business Logic",
            font_size=40,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Add comparison infographic
        img_buf = self.create_code_comparison()
        slide.shapes.add_picture(
            img_buf, Inches(0.5), Inches(1.2), width=Inches(12.333)
        )

    def create_architecture_slide(self):
        """Slide 7: Architecture Enforcement"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Guardrails, Not Gates",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Create flow diagram
        img_buf = BytesIO()
        fig, ax = plt.subplots(figsize=(12, 5))
        ax.set_xlim(0, 12)
        ax.set_ylim(0, 5)
        ax.axis("off")

        # Bronze box
        bronze_box = FancyBboxPatch(
            (0.5, 2),
            2.5,
            1.5,
            boxstyle="round,pad=0.1",
            facecolor="#CD7F32",
            edgecolor="#8B4513",
            linewidth=2,
        )
        ax.add_patch(bronze_box)
        ax.text(
            1.75,
            2.75,
            "Bronze",
            ha="center",
            va="center",
            fontsize=14,
            fontweight="bold",
            color="white",
        )

        # Invalid arrow (blocked)
        ax.annotate(
            "",
            xy=(4, 2.75),
            xytext=(3.2, 2.75),
            arrowprops=dict(arrowstyle="->", color="#DC2626", lw=3),
        )
        ax.text(3.6, 2.3, "❌", fontsize=20, ha="center")

        # Gold box (blocked)
        gold_box_invalid = FancyBboxPatch(
            (4, 2),
            2.5,
            1.5,
            boxstyle="round,pad=0.1",
            facecolor="#FFD700",
            edgecolor="#DC2626",
            linewidth=3,
            linestyle="--",
        )
        ax.add_patch(gold_box_invalid)
        ax.text(
            5.25,
            2.75,
            "Gold\n(BLOCKED)",
            ha="center",
            va="center",
            fontsize=12,
            fontweight="bold",
        )

        # Valid flow
        ax.annotate(
            "",
            xy=(3.2, 1.5),
            xytext=(3.2, 1.9),
            arrowprops=dict(arrowstyle="->", color="#10B981", lw=3),
        )

        # Silver box
        silver_box = FancyBboxPatch(
            (0.5, 0.2),
            2.5,
            1.5,
            boxstyle="round,pad=0.1",
            facecolor="#C0C0C0",
            edgecolor="#6B7280",
            linewidth=2,
        )
        ax.add_patch(silver_box)
        ax.text(
            1.75,
            0.95,
            "Silver",
            ha="center",
            va="center",
            fontsize=14,
            fontweight="bold",
        )

        # Valid arrow
        ax.annotate(
            "",
            xy=(4, 0.95),
            xytext=(3.2, 0.95),
            arrowprops=dict(arrowstyle="->", color="#10B981", lw=3),
        )
        ax.text(3.6, 0.5, "✅", fontsize=20, ha="center")

        # Gold box (valid)
        gold_box_valid = FancyBboxPatch(
            (4, 0.2),
            2.5,
            1.5,
            boxstyle="round,pad=0.1",
            facecolor="#FFD700",
            edgecolor="#10B981",
            linewidth=3,
        )
        ax.add_patch(gold_box_valid)
        ax.text(
            5.25,
            0.95,
            "Gold\n(ALLOWED)",
            ha="center",
            va="center",
            fontsize=12,
            fontweight="bold",
        )

        # Key points
        points = [
            "✓ Cannot skip layers",
            "✓ Automatic dependency tracking",
            "✓ Type safety between stages",
            "✓ Built-in validation patterns",
        ]

        y = 4
        for point in points:
            ax.text(7.5, y, point, fontsize=12, va="center")
            y -= 0.6

        plt.savefig(
            img_buf, format="png", dpi=150, bbox_inches="tight", facecolor="white"
        )
        img_buf.seek(0)
        plt.close()

        slide.shapes.add_picture(
            img_buf, Inches(0.5), Inches(1.5), width=Inches(12.333)
        )

    def create_impact_slide(self):
        """Slide 8: Real-World Impact"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Results That Matter",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Add metrics chart
        img_buf = self.create_metrics_chart()
        slide.shapes.add_picture(img_buf, Inches(1), Inches(1.5), width=Inches(11.333))

    def create_ecosystem_slide(self):
        """Slide 9: Ecosystem Integration"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Works With Your Stack",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Integration categories
        integrations = [
            ("Languages & Frameworks", ["🐍 Python", "⚡ Polars", "🐼 Pandas"]),
            ("Orchestration", ["🔄 Dagster", "☁️ Cloud-native", "🚀 Scalable"]),
            ("Data Sources", ["🐘 PostgreSQL", "📊 MySQL", "🔌 REST APIs", "📁 Files"]),
            ("Destinations", ["☁️ AWS", "☁️ Azure", "☁️ GCP", "📈 BI Tools"]),
        ]

        y_pos = 1.8
        for category, items in integrations:
            self.add_text_box(
                slide,
                0.5,
                y_pos,
                3.5,
                0.5,
                category,
                font_size=16,
                bold=True,
                font_color=COLORS["dagster_purple"],
            )
            items_text = " • ".join(items)
            self.add_text_box(
                slide,
                4,
                y_pos,
                8.5,
                0.5,
                items_text,
                font_size=14,
                font_color=COLORS["text_gray"],
            )
            y_pos += 0.8

        # Hub diagram
        img_buf = BytesIO()
        fig, ax = plt.subplots(figsize=(6, 6))
        ax.set_xlim(0, 6)
        ax.set_ylim(0, 6)
        ax.axis("off")

        # Center hub - Polster
        center = Circle(
            (3, 3), 0.8, facecolor="#7C3AED", edgecolor="#6D28D9", linewidth=3
        )
        ax.add_patch(center)
        ax.text(
            3,
            3,
            "Polster",
            ha="center",
            va="center",
            fontsize=12,
            fontweight="bold",
            color="white",
        )

        # Surrounding nodes
        nodes = [
            (3, 5, "Python"),
            (5, 4, "Dagster"),
            (5, 2, "Polars"),
            (3, 1, "Cloud"),
            (1, 2, "SQL"),
            (1, 4, "APIs"),
        ]

        for x, y, label in nodes:
            node = Circle(
                (x, y), 0.5, facecolor="#2874A6", edgecolor="#1E40AF", linewidth=2
            )
            ax.add_patch(node)
            ax.text(
                x,
                y,
                label,
                ha="center",
                va="center",
                fontsize=9,
                color="white",
                fontweight="bold",
            )
            # Connection line
            ax.plot([3, x], [3, y], "o-", color="#CBD5E1", linewidth=2, alpha=0.6)

        plt.savefig(
            img_buf, format="png", dpi=150, bbox_inches="tight", facecolor="white"
        )
        img_buf.seek(0)
        plt.close()

        slide.shapes.add_picture(img_buf, Inches(8), Inches(2.5), width=Inches(4.5))

    def create_use_cases_slide(self):
        """Slide 10: Use Cases"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Who Is Polster For?",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Three personas
        personas = [
            (
                "Data Engineer",
                "Finally, a standardized way to build pipelines that my team can actually maintain.",
                "💻",
            ),
            (
                "Analytics Manager",
                "My team ships faster with fewer bugs. New hires get productive in days, not weeks.",
                "📊",
            ),
            (
                "Startup CTO",
                "Enterprise-grade data infrastructure without the enterprise timeline or budget.",
                "🚀",
            ),
        ]

        x_positions = [0.5, 4.5, 8.5]
        for i, (role, quote, icon) in enumerate(personas):
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_positions[i]),
                Inches(2),
                Inches(3.8),
                Inches(4.5),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(248, 250, 252)
            card.line.color.rgb = (
                COLORS["polars_blue"] if i % 2 == 0 else COLORS["dagster_purple"]
            )
            card.line.width = Pt(2)

            # Icon
            self.add_text_box(
                slide, x_positions[i] + 1.4, 2.2, 1, 0.8, icon, font_size=40
            )

            # Role
            self.add_text_box(
                slide,
                x_positions[i] + 0.2,
                3.1,
                3.4,
                0.6,
                role,
                font_size=20,
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )

            # Quote
            self.add_text_box(
                slide,
                x_positions[i] + 0.2,
                3.8,
                3.4,
                2.5,
                f'"{quote}"',
                font_size=14,
                alignment=PP_ALIGN.CENTER,
                font_color=COLORS["text_gray"],
            )

    def create_getting_started_slide(self):
        """Slide 11: Getting Started"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self.add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            "Start Building Today",
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # 3-step process
        steps = [
            ("1", "pip install polster", "Install from PyPI"),
            ("2", "polster init my_project", "Create your first pipeline"),
            ("3", "python run_polster.py --ui", "Launch the UI"),
        ]

        y_pos = 2
        colors = [COLORS["bronze"], COLORS["silver"], COLORS["gold"]]

        for i, (num, cmd, desc) in enumerate(steps):
            # Step number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.8), Inches(0.8)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors[i]
            circle.line.color.rgb = RGBColor(0, 0, 0)

            self.add_text_box(
                slide,
                1,
                y_pos + 0.15,
                0.8,
                0.5,
                num,
                font_size=24,
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )

            # Command
            cmd_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2),
                Inches(y_pos),
                Inches(6),
                Inches(0.8),
            )
            cmd_box.fill.solid()
            cmd_box.fill.fore_color.rgb = RGBColor(15, 23, 42)

            self.add_text_box(
                slide,
                2.2,
                y_pos + 0.2,
                5.6,
                0.4,
                f"$ {cmd}",
                font_size=18,
                font_color=COLORS["white"],
                font_name="Courier New",
            )

            # Description
            self.add_text_box(
                slide,
                8.5,
                y_pos + 0.2,
                4,
                0.4,
                desc,
                font_size=16,
                font_color=COLORS["text_gray"],
            )

            y_pos += 1.5

        # QR Code placeholder and links
        link_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(5.8),
            Inches(12.333),
            Inches(1.2),
        )
        link_box.fill.solid()
        link_box.fill.fore_color.rgb = RGBColor(224, 231, 255)
        link_box.line.color.rgb = COLORS["dagster_purple"]

        self.add_text_box(
            slide,
            0.7,
            6,
            12,
            0.8,
            "📦 pip install polster  •  🐙 github.com/sultanaltair96/polster-cli  •  📚 polster.rtfd.io",
            font_size=16,
            alignment=PP_ALIGN.CENTER,
        )

    def create_closing_slide(self):
        """Slide 12: Closing"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Large logo text
        self.add_text_box(
            slide,
            0.5,
            2.5,
            12.333,
            1.5,
            "Polster",
            font_size=80,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Tagline
        self.add_text_box(
            slide,
            0.5,
            4,
            12.333,
            1,
            "Architecture that scales with your data",
            font_size=28,
            alignment=PP_ALIGN.CENTER,
            font_color=COLORS["text_gray"],
        )

        # Contact info
        self.add_text_box(
            slide,
            0.5,
            5.5,
            12.333,
            0.5,
            "🚀 pip install polster  |  💻 github.com/sultanaltair96/polster-cli",
            font_size=18,
            alignment=PP_ALIGN.CENTER,
        )

        self.add_text_box(
            slide,
            0.5,
            6.2,
            12.333,
            0.5,
            "Questions? Let's talk!",
            font_size=20,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_color=COLORS["dagster_purple"],
        )

    def generate_presentation(self, output_path="polster_presentation.pptx"):
        """Generate the complete presentation"""
        print("🎨 Generating Polster Presentation...")

        # Create all slides
        self.create_title_slide()
        print("✓ Slide 1: Title")

        self.create_problem_slide()
        print("✓ Slide 2: The Problem")

        self.create_medallion_slide()
        print("✓ Slide 3: Medallion Architecture")

        self.create_polster_intro_slide()
        print("✓ Slide 4: Meet Polster")

        self.create_demo_slide()
        print("✓ Slide 5: Live Demo")

        self.create_code_slide()
        print("✓ Slide 6: Code Comparison")

        self.create_architecture_slide()
        print("✓ Slide 7: Architecture Enforcement")

        self.create_impact_slide()
        print("✓ Slide 8: Impact Metrics")

        self.create_ecosystem_slide()
        print("✓ Slide 9: Ecosystem")

        self.create_use_cases_slide()
        print("✓ Slide 10: Use Cases")

        self.create_getting_started_slide()
        print("✓ Slide 11: Getting Started")

        self.create_closing_slide()
        print("✓ Slide 12: Closing")

        # Save presentation
        self.prs.save(output_path)
        print(f"\n✅ Presentation saved: {output_path}")
        print(f"📊 Total slides: {len(self.prs.slides)}")
        print(f"🎨 Theme: White with Polars Blue & Dagster Purple")


def main():
    """Main function to generate the presentation"""
    # Create presentation
    presentation = PolsterPresentation()
    presentation.generate_presentation("polster_presentation.pptx")

    # Also generate PDF version info
    print("\n💡 Tips for presenting:")
    print("   • Duration: 15-20 minutes")
    print("   • Interactive demo: Slides 4-5")
    print("   • Q&A: After slide 12")
    print("\n📄 To convert to PDF:")
    print("   • Open in PowerPoint → File → Export → Create PDF/XPS")
    print(
        "   • Or use: LibreOffice --headless --convert-to pdf polster_presentation.pptx"
    )


if __name__ == "__main__":
    main()
